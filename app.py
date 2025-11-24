import os
import json
import uuid
import numpy as np
import streamlit as st
from openai import OpenAI
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# ---------------- CONFIG ----------------
DATA_ROOT = "data"
INDEX_DIR = "index"
CHUNKS_PATH = os.path.join(INDEX_DIR, "chunks.jsonl")
EMB_PATH = os.path.join(INDEX_DIR, "embeddings.npy")
EMBED_MODEL = "text-embedding-3-small"
GPT_MODEL = "gpt-4.1-mini"   # or gpt-4.1, gpt-4o-mini, etc.


os.makedirs(DATA_ROOT, exist_ok=True)
os.makedirs(INDEX_DIR, exist_ok=True)

client = OpenAI()

# In-memory cache
_embeddings_cache = None
_chunks_cache = None

# ---------------- INDEX UTILS ----------------

def load_index():
    """Load chunks + embeddings into memory (if they exist)."""
    global _embeddings_cache, _chunks_cache

    if not os.path.exists(CHUNKS_PATH) or not os.path.exists(EMB_PATH):
        _embeddings_cache = np.zeros((0, 1536), dtype=np.float32)
        _chunks_cache = []
        return

    chunks = []
    with open(CHUNKS_PATH, "r", encoding="utf-8") as f:
        for line in f:
            chunks.append(json.loads(line))
    _chunks_cache = chunks

    _embeddings_cache = np.load(EMB_PATH)


def save_index():
    """Write current chunks + embeddings to disk."""
    global _embeddings_cache, _chunks_cache

    with open(CHUNKS_PATH, "w", encoding="utf-8") as f:
        for ch in _chunks_cache:
            f.write(json.dumps(ch, ensure_ascii=False) + "\n")

    np.save(EMB_PATH, _embeddings_cache)


def embed_text(text: str) -> np.ndarray:
    resp = client.embeddings.create(
        model=EMBED_MODEL,
        input=text
    )
    return np.array(resp.data[0].embedding, dtype=np.float32)

# ---------------- TEXT EXTRACTION ----------------

def extract_text_from_pdf(filepath: str) -> str:
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def extract_text_from_docx(filepath: str) -> str:
    doc = DocxDocument(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_xlsx(filepath: str) -> str:
    # Simple: concatenate all sheets into a CSV-like text
    text_parts = []
    xls = pd.ExcelFile(filepath)
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        text_parts.append(f"=== SHEET: {sheet_name} ===")
        text_parts.append(df.to_csv(index=False))
    return "\n".join(text_parts)


def extract_text_from_ppt(filepath: str) -> str:
    prs = Presentation(filepath)
    text_parts = []
    for i, slide in enumerate(prs.slides):
        text_parts.append(f"=== SLIDE {i+1} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
    return "\n".join(text_parts)


def extract_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == ".pdf":
            return extract_text_from_pdf(filepath)
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif ext == ".docx":
            return extract_text_from_docx(filepath)
        elif ext in [".xlsx", ".xls"]:
            return extract_text_from_xlsx(filepath)
        elif ext in [".pptx", ".ppt"]:
            return extract_text_from_ppt(filepath)
        else:
            # Unsupported for now
            return ""
    except Exception as e:
        # Avoid crashing the whole app just because one file fails
        return f"[ERROR extracting text from {os.path.basename(filepath)}: {e}]"

# ---------------- CHUNKING + SIM ----------------

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200):
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks


def cosine_sim(a: np.ndarray, b: np.ndarray) -> np.ndarray:
    a_norm = a / (np.linalg.norm(a) + 1e-8)
    b_norm = b / (np.linalg.norm(b, axis=1, keepdims=True) + 1e-8)
    return np.dot(b_norm, a_norm)

# ---------------- HELPERS ----------------

def available_projects():
    if not os.path.exists(DATA_ROOT):
        return ["All"]
    return ["All"] + [
        name for name in os.listdir(DATA_ROOT)
        if os.path.isdir(os.path.join(DATA_ROOT, name))
    ]


def get_documents_overview():
    """
    Build a unique list of documents from the chunks index.
    Each row: project, folder, filename, filepath.
    """
    global _chunks_cache
    if _chunks_cache is None:
        load_index()

    docs_seen = {}
    for ch in _chunks_cache:
        key = ch.get("filepath")
        if not key:
            continue
        if key not in docs_seen:
            docs_seen[key] = {
                "project": ch.get("project", ""),
                "folder": ch.get("folder", ""),
                "filename": ch.get("filename", ""),
                "filepath": ch.get("filepath", ""),
            }
    if not docs_seen:
        return pd.DataFrame(columns=["project", "folder", "filename", "filepath"])
    return pd.DataFrame(list(docs_seen.values()))


# ---------------- UPLOAD LOGIC ----------------

def handle_upload(files, project, folder):
    global _embeddings_cache, _chunks_cache

    if not files:
        return "No files uploaded."

    if _embeddings_cache is None or _chunks_cache is None:
        load_index()

    project = project.strip() or "General"
    folder = folder.strip() or "misc"

    target_dir = os.path.join(DATA_ROOT, project, folder)
    os.makedirs(target_dir, exist_ok=True)

    # First pass: save files and extract chunks so we can show accurate progress
    file_chunks = []
    added_docs = []
    total_chunks = 0

    for f in files:
        filename = f.name
        save_path = os.path.join(target_dir, filename)

        with open(save_path, "wb") as out:
            out.write(f.getbuffer())

        text = extract_text(save_path)
        if not text.strip():
            # still consider saved but no chunks
            added_docs.append(f"{filename} → {project}/{folder} (no text extracted)")
            continue

        chunks = chunk_text(text)
        file_chunks.append((save_path, filename, chunks))
        total_chunks += len(chunks)
        added_docs.append(f"{filename} → {project}/{folder}")

    if total_chunks == 0:
        # Nothing to index
        save_index()
        return "Files saved but no text was extracted (check formats)."

    # Progress UI
    progress = st.progress(0)
    status = st.empty()
    processed = 0

    for save_path, filename, chunks in file_chunks:
        status.text(f"Indexing {filename}...")
        for chunk in chunks:
            emb = embed_text(chunk)

            meta = {
                "id": str(uuid.uuid4()),
                "project": project,
                "folder": folder,
                "filename": filename,
                "filepath": save_path,
                "text": chunk[:4000]   # keep per-chunk text
            }

            if _embeddings_cache.shape[0] == 0:
                _embeddings_cache = emb.reshape(1, -1)
            else:
                _embeddings_cache = np.vstack([_embeddings_cache, emb])

            _chunks_cache.append(meta)

            processed += 1
            progress.progress(int(processed / total_chunks * 100))

    # done
    status.empty()
    progress.empty()
    save_index()
    return "Indexed:\n" + "\n".join(added_docs)


def scan_and_index_disk():
    """
    Walk DATA_ROOT and (re)build the index from files on disk with progress.
    This does a full rebuild of _embeddings_cache and _chunks_cache.
    """
    global _embeddings_cache, _chunks_cache

    allowed_exts = {".pdf", ".docx", ".txt", ".xlsx", ".xls", ".ppt", ".pptx"}
    files = []
    for root, _, filenames in os.walk(DATA_ROOT):
        for fn in filenames:
            if os.path.splitext(fn)[1].lower() in allowed_exts:
                files.append(os.path.join(root, fn))

    if not files:
        return "No files found under data/ to index."

    # First pass: extract chunks for each file to compute total work
    files_with_chunks = []
    total_chunks = 0
    for fp in sorted(files):
        try:
            text = extract_text(fp)
            if not text.strip():
                continue
            chunks = chunk_text(text)
            if not chunks:
                continue
            rel = os.path.relpath(fp, DATA_ROOT).split(os.sep)
            project = rel[0] if len(rel) > 0 else "General"
            folder = rel[1] if len(rel) > 1 else "misc"
            files_with_chunks.append((fp, os.path.basename(fp), project, folder, chunks))
            total_chunks += len(chunks)
        except Exception as e:
            files_with_chunks.append((fp, os.path.basename(fp), "ERROR", "ERROR", [f"[ERROR extracting: {e}]"]))
            total_chunks += 1

    if total_chunks == 0:
        return "No extractable text found in data/."

    # Reset caches
    _embeddings_cache = np.zeros((0, 1536), dtype=np.float32)
    _chunks_cache = []

    progress = st.progress(0)
    status = st.empty()
    processed = 0
    indexed = []

    for fp, basename, project, folder, chunks in files_with_chunks:
        status.text(f"Indexing {basename}...")
        mtime = os.path.getmtime(fp) if os.path.exists(fp) else None
        for chunk in chunks:
            emb = embed_text(chunk)
            meta = {
                "id": str(uuid.uuid4()),
                "project": project,
                "folder": folder,
                "filename": basename,
                "filepath": fp,
                "mtime": mtime,
                "text": chunk[:4000],
            }
            if _embeddings_cache.shape[0] == 0:
                _embeddings_cache = emb.reshape(1, -1)
            else:
                _embeddings_cache = np.vstack([_embeddings_cache, emb])
            _chunks_cache.append(meta)

            processed += 1
            progress.progress(int(processed / total_chunks * 100))

        indexed.append(f"{basename} → {project}/{folder}")

    status.empty()
    progress.empty()
    save_index()
    return "Reindexed:\n" + "\n".join(indexed)

# ---------------- CHAT / RAG LOGIC ----------------

def rag_answer(
    question,
    project_filter,
    sim_threshold: float = 0.60,
    max_files: int = 5,
    max_sources: int = 3,
):
    """
    Two-stage retrieval strategy:
      1) Try direct semantic retrieval (strict).
      2) If no candidates, rewrite the question into a search query and retry (more recall).
      3) If still no candidates, return a generative fallback clearly marked as UNGROUNDED.
    """
    global _embeddings_cache, _chunks_cache

    if _embeddings_cache is None or _chunks_cache is None:
        load_index()

    if _embeddings_cache.shape[0] == 0:
        return "No documents indexed yet."

    def retrieve_for_query(q, threshold, max_files_local):
        q_emb_local = embed_text(q)
        sims_local = cosine_sim(q_emb_local, _embeddings_cache).flatten()
        sorted_idxs_local = np.argsort(-sims_local)

        candidates_local = []
        for idx in sorted_idxs_local:
            sim = float(sims_local[int(idx)])
            meta = _chunks_cache[int(idx)]

            if project_filter != "All" and meta.get("project") != project_filter:
                continue

            if sim < threshold:
                continue

            candidates_local.append({"meta": meta, "sim": sim})
            if len(candidates_local) >= max_files_local:
                break

        return candidates_local

    # 1) Try strict retrieval
    candidates = retrieve_for_query(question, sim_threshold, max_files)

    # 2) If nothing, ask the model to rewrite the query and retry (helps indirect/phrased questions)
    if not candidates:
        try:
            rewrite_prompt = (
                "Rewrite the user's latest question into a concise search query (1-2 sentences) "
                "for QA/QE documents. Return only the rewritten query.\n\n"
                f"Latest user question: {question}\n\nRewrite:"
            )
            resp = client.responses.create(model=GPT_MODEL, input=rewrite_prompt)
            try:
                rewritten = resp.output[0].content[0].text.strip()
            except Exception:
                rewritten = str(resp).strip()
            if rewritten:
                # retry with rewritten query using a bit lower threshold and larger max_files
                candidates = retrieve_for_query(rewritten, max(sim_threshold * 0.8, 0.30), max_files * 3)
        except Exception:
            candidates = []

    # 3) If still nothing, generative fallback (clearly labeled)
    if not candidates:
        try:
            gen_prompt = (
                "You are a Quality Engineering assistant. The documents were searched but no sufficiently "
                "relevant passages were found. Answer the question using your general knowledge, but prefix the "
                "answer with the text '[UNGROUNDED]' so the user knows this is not supported by indexed docs.\n\n"
                f"Question: {question}\n\nAnswer:"
            )
            resp2 = client.responses.create(model=GPT_MODEL, input=gen_prompt)
            try:
                gen_answer = resp2.output[0].content[0].text
            except Exception:
                gen_answer = str(resp2)
            return gen_answer + "\n\n**Note:** [UNGROUNDED] — no verified sources found in the indexed documents. Answers are from general AI knowledge and may be inaccurate."
        except Exception as e:
            return f"I don't know — no relevant documents found and generative fallback failed: {e}"

    # Build prompt from retrieved candidates (same safe behavior you already have)
    context = "\n\n".join(
        f"[{c['meta']['filename']} | {c['meta']['project']} | {c['meta']['folder']}] (sim={c['sim']:.3f})\n{c['meta']['text']}"
        for c in candidates
    )

    prompt = f"""
You are a Quality Engineering Document Assistant.
Answer ONLY using the CONTEXT below. If the answer is not present in the CONTEXT, respond exactly: "I don't know — not enough information in the documents."
Cite the document filename(s) you used at the end.

CONTEXT:
{context}

QUESTION:
{question}
"""

    resp = client.responses.create(model=GPT_MODEL, input=prompt)
    try:
        answer = resp.output[0].content[0].text
    except Exception:
        answer = str(resp)

    # Deduplicate by filepath and keep top `max_sources` by similarity
    seen = set()
    deduped = []
    for c in candidates:
        path = c['meta'].get('filepath') or f"{c['meta'].get('project')}|{c['meta'].get('filename')}"
        if path in seen:
            continue
        seen.add(path)
        deduped.append(c)
        if len(deduped) >= max_sources:
            break

    sources = "\n".join(
        f"- {d['meta']['filename']} (project: {d['meta']['project']}, folder: {d['meta']['folder']}, sim={d['sim']:.3f})"
        for d in deduped
    )

    return answer + "\n\n**Sources:**\n" + sources

# ---------------- STREAMLIT UI ----------------

st.set_page_config(page_title="QE Local RAG Assistant", layout="wide")
st.title("(QEAi) Quality Engineering AI Assistant ")

tab_chat, tab_upload, tab_browse = st.tabs(["💬 Chat", "📤 Upload & Index", "📂 Browse Documents"])

# ----- CHAT TAB -----
with tab_chat:
    st.subheader("Ask questions about your QE documents")

    project_filter = st.selectbox(
        "Filter by project:",
        available_projects(),
        key="chat_project_filter"
    )

    # Preset strictness dropdown (replaces slider / checkbox)
    preset_options = ["Strict (safer)", "Balanced", "Lenient (more answers, less reliable)"]
    preset = st.selectbox(
        "Answer strictness preset:",
        preset_options,
        index=0,
        help="Strict: higher similarity threshold to reduce hallucinations. Balanced: moderate. Lenient: lower threshold to increase recall."
    )

    # map presets to similarity thresholds
    preset_thresholds = {
        "Strict (safer)": 0.60,
        "Balanced": 0.50,
        "Lenient (more answers, less reliable)": 0.40
    }
    sim_threshold = preset_thresholds.get(preset, 0.60)

    # show a small notice when not strict
    if preset != "Strict (safer)":
        st.warning("Using a less strict preset — responses may be less reliable. Verify quoted snippets and sources before trusting the answer.")

    st.markdown("### Question")
    question = st.text_area(
        "Question:",
        height=120,
        placeholder="e.g., What is the cable bend requirement for a Fiber cable?",
    )
    if st.button("Send", type="primary"):
        question_text = (question or "").strip()
        if not question_text:
            st.warning("Please enter a question.")
        else:
            with st.spinner("Thinking..."):
                answer = rag_answer(
                    question_text,
                    project_filter,
                    sim_threshold=sim_threshold,
                )
            st.markdown("### Answer")
            st.markdown(answer)

# ----- UPLOAD TAB -----
with tab_upload:
    st.subheader("Upload and index documents")

    uploaded_files = st.file_uploader(
        "Upload files (PDF, DOCX, XLSX, PPTX, TXT):",
        accept_multiple_files=True,
        type=["pdf", "docx", "xlsx", "xls", "ppt", "pptx", "txt"]
    )

    # project / folder pickers: show existing values and allow new entry
    col1, col2 = st.columns(2)
    with col1:
        existing_projects = [
            name for name in os.listdir(DATA_ROOT)
            if os.path.isdir(os.path.join(DATA_ROOT, name))
        ] if os.path.exists(DATA_ROOT) else []
        project_choice_options = ["<New project>"] + sorted(existing_projects) if existing_projects else ["<New project>"]
        project_choice = st.selectbox("Project (choose existing or New):", project_choice_options, index=0)

        if project_choice == "<New project>":
            project = st.text_input("New project name", value="")
        else:
            project = project_choice

    with col2:
        # list folders for selected existing project
        folder_options = []
        if project_choice != "<New project>" and project_choice:
            proj_path = os.path.join(DATA_ROOT, project_choice)
            if os.path.exists(proj_path):
                folder_options = [
                    name for name in os.listdir(proj_path)
                    if os.path.isdir(os.path.join(proj_path, name))
                ]
        folder_choice_options = ["<New folder>"] + sorted(folder_options) if folder_options else ["<New folder>"]
        folder_choice = st.selectbox("Folder (choose existing or New):", folder_choice_options, index=0)

        if folder_choice == "<New folder>":
            folder = st.text_input("New folder name", value="")
        else:
            folder = folder_choice

    # (doc_type and tags removed as requested)

    if st.button("Upload & Index", type="primary"):
        if not project or not project.strip():
            st.warning("Please provide a project name.")
        elif not folder or not folder.strip():
            st.warning("Please provide a folder name.")
        else:
            with st.spinner("Processing and indexing..."):
                message = handle_upload(uploaded_files, project, folder)
            st.success(message)

    st.divider()
    st.markdown("Or: place files directly into the data/ folder and click:")
    if st.button("Scan data/ folder and (re)build index"):
        with st.spinner("Scanning data/ and indexing files (may take a while)..."):
            msg = scan_and_index_disk()
        st.success(msg)

# ----- BROWSE TAB -----
with tab_browse:
    # ensure index is loaded and build the documents DataFrame
    if _chunks_cache is None or _embeddings_cache is None:
        load_index()
    df_docs = get_documents_overview()

    colp, colf = st.columns(2)
    with colp:
        project_options = ["All"] + sorted(list(set(df_docs["project"]))) if not df_docs.empty else ["All"]
        proj_sel = st.selectbox(
            "Filter by project:",
            project_options,
            key="browse_project_filter"
        )
    with colf:
        folder_options = ["All"] + sorted(list(set(df_docs["folder"]))) if not df_docs.empty else ["All"]
        folder_sel = st.selectbox(
            "Filter by folder:",
            folder_options,
            key="browse_folder_filter"
        )

    filtered = df_docs.copy()
    if proj_sel != "All":
        filtered = filtered[filtered["project"] == proj_sel]
    if folder_sel != "All":
        filtered = filtered[filtered["folder"] == folder_sel]

    st.dataframe(
        filtered[["project", "folder", "filename"]],
        width="stretch"
    )
    st.caption("Note: actual files are stored under the local 'data/' directory on disk.")
